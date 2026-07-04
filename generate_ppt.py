"""
Generate a new PPT for ISO 11820 建筑材料不燃性试验仿真系统
Following the same structure as the original template PPT.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# Constants
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Color scheme - professional blue theme
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
MEDIUM_BLUE = RGBColor(0x2B, 0x57, 0x9A)
LIGHT_BLUE = RGBColor(0x3D, 0x7E, 0xC7)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
BG_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
YELLOW_BG = RGBColor(0xFF, 0xF3, 0xCD)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout as base
blank_layout = prs.slide_layouts[6]  # blank layout


def add_bg_rect(slide, left, top, width, height, color):
    """Add a solid color rectangle as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 font_color=BLACK, bold=False, font_name='Microsoft YaHei',
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with single style."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=Pt(12),
                       default_color=DARK_GRAY, default_name='Microsoft YaHei'):
    """Add a text box with multiple styled lines. lines is a list of (text, size, color, bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, False
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else False

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = default_name
        p.space_after = Pt(4)
    return txBox


def add_section_number(slide, top_offset=Inches(0.45)):
    """Add section number styling - a large colored rectangle behind the number."""
    # Just return the standard position
    return Inches(0.4), top_offset


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Dark blue background
add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)

# Decorative accent bar at bottom
add_bg_rect(slide, 0, Inches(4.8), SLIDE_W, Inches(0.05), ACCENT_BLUE)

# Main title
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.2),
             '建筑材料不燃性试验仿真系统\n设计与实现',
             font_size=Pt(36), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle line
add_bg_rect(slide, Inches(3.5), Inches(2.55), Inches(3), Inches(0.03), ACCENT_BLUE)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.6),
             '——毕业设计答辩',
             font_size=Pt(22), font_color=RGBColor(0xCC, 0xDD, 0xFF), bold=False,
             alignment=PP_ALIGN.CENTER)

# Bottom info
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.5),
             '基于 .NET 8 与 WinForms 的 ISO 11820 标准仿真测控平台',
             font_size=Pt(14), font_color=RGBColor(0xAA, 0xBB, 0xDD), bold=False,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(blank_layout)

add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
# Top bar
add_bg_rect(slide, 0, 0, SLIDE_W, Inches(0.08), DARK_BLUE)

# Title
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
             '目录', font_size=Pt(30), font_color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.95), Inches(8.4), Inches(0.3),
             '· CONTENTS ·', font_size=Pt(12), font_color=MED_GRAY)

# TOC items
toc_items = [
    ('01', '研究背景与意义', 'BACKGROUND AND SIGNIFICANCE'),
    ('02', '系统需求分析与技术选型', 'REQUIREMENTS & TECHNOLOGY SELECTION'),
    ('03', '系统总体设计', 'SYSTEM DESIGN'),
    ('04', '系统详细设计与实现', 'DETAILED DESIGN & IMPLEMENTATION'),
]

for idx, (num, cn, en) in enumerate(toc_items):
    y = Inches(1.6) + Inches(0.85) * idx

    # Number circle/box
    num_box = add_text_box(slide, Inches(1.0), y, Inches(0.7), Inches(0.6),
                           num, font_size=Pt(28), font_color=DARK_BLUE, bold=True,
                           alignment=PP_ALIGN.CENTER)

    # Separator line
    add_bg_rect(slide, Inches(1.9), y + Inches(0.15), Inches(0.04), Inches(0.3), ACCENT_BLUE)

    # Chinese title
    add_text_box(slide, Inches(2.2), y + Inches(0.0), Inches(5), Inches(0.35),
                 cn, font_size=Pt(20), font_color=DARK_BLUE, bold=True)

    # English subtitle
    add_text_box(slide, Inches(2.2), y + Inches(0.35), Inches(5), Inches(0.25),
                 en, font_size=Pt(10), font_color=MED_GRAY)

# Bottom line
add_bg_rect(slide, 0, Inches(5.3), SLIDE_W, Inches(0.05), DARK_BLUE)


# ============================================================
# Helper for section divider slides
# ============================================================
def make_section_slide(prs, section_num, cn_title, en_title):
    slide = prs.slides.add_slide(blank_layout)

    # Outer border
    add_bg_rect(slide, Inches(0.15), Inches(0.05), Inches(9.7), Inches(5.45), WHITE)
    add_bg_rect(slide, Inches(0.3), Inches(0.2), Inches(9.4), Inches(5.15), LIGHT_GRAY)

    # Large number background
    add_bg_rect(slide, Inches(0.5), Inches(0.4), Inches(3.9), Inches(1.2), DARK_BLUE)

    # Number text
    add_text_box(slide, Inches(0.4), Inches(0.35), Inches(4.1), Inches(1.3),
                 section_num, font_size=Pt(64), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Horizontal line
    add_bg_rect(slide, Inches(0.5), Inches(2.0), Inches(0.9), Inches(0.006), ACCENT_BLUE)

    # CN Title background
    add_bg_rect(slide, Inches(0.5), Inches(2.2), Inches(4.7), Inches(0.5), DARK_BLUE)
    add_text_box(slide, Inches(0.4), Inches(2.15), Inches(4.9), Inches(0.55),
                 cn_title, font_size=Pt(28), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # EN Title
    add_bg_rect(slide, Inches(0.5), Inches(2.75), Inches(4.7), Inches(0.25), ACCENT_BLUE)
    add_text_box(slide, Inches(0.4), Inches(2.72), Inches(4.9), Inches(0.3),
                 en_title, font_size=Pt(12), font_color=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Vertical line on the right
    add_bg_rect(slide, Inches(8.6), Inches(0.6), Inches(0.006), Inches(3.9), MED_GRAY)

    return slide


def make_content_slide(prs, title_text, page_num=None):
    """Create a standard content slide with header bar."""
    slide = prs.slides.add_slide(blank_layout)

    # White background
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Top bar
    add_bg_rect(slide, 0, 0, SLIDE_W, Inches(0.06), DARK_BLUE)

    # Title area with blue accent
    add_bg_rect(slide, Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.55), DARK_BLUE)
    add_text_box(slide, Inches(0.5), Inches(0.22), Inches(9), Inches(0.5),
                 title_text, font_size=Pt(18), font_color=WHITE, bold=True)

    # Page number
    if page_num:
        add_text_box(slide, Inches(8.8), Inches(5.2), Inches(1), Inches(0.3),
                     f'第{page_num}页', font_size=Pt(9), font_color=MED_GRAY,
                     alignment=PP_ALIGN.RIGHT)

    # Bottom line
    add_bg_rect(slide, 0, Inches(5.3), SLIDE_W, Inches(0.04), DARK_BLUE)

    return slide


# ============================================================
# SLIDE 3: Section 01 Divider
# ============================================================
make_section_slide(prs, '01', '研究背景与意义', 'BACKGROUND AND SIGNIFICANCE')


# ============================================================
# SLIDE 4: Research Background
# ============================================================
slide = make_content_slide(prs, '1. 研究背景及意义 · 研究背景', 4)

add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0), [
    ('建筑材料的防火安全性能是构筑社会公共安全防线的基石，直接关系到建筑结构在火灾中的稳定性、人员疏散安全以及财产保护。不燃性作为评价建筑材料及制品防火等级最基础和最重要的性能指标之一，其测试结果的准确性、可靠性和可复现性具有法规强制性意义。', Pt(14), DARK_GRAY, False),
    ('', Pt(8), DARK_GRAY, False),
    ('ISO 11820-2020《建筑材料不燃性试验方法》是国际公认的评估建筑材料防火性能的核心标准。该标准严格规定了试验必须在750°C ± 5°C的恒定炉温条件下进行，记录30至60分钟的温度数据，并依据温升（ΔT ≤ 50°C）和失重率（≤ 50%）等关键指标对材料的不燃性能作出判定。', Pt(14), DARK_GRAY, False),
    ('', Pt(8), DARK_GRAY, False),
    ('然而，完整的试验硬件设备（加热炉、PID控制器、传感器组、Modbus通信模块等）造价昂贵、占地大且维护复杂，不利于教学环境中的大规模部署。因此，开发一套高保真的仿真系统——以软件方式模拟完整的试验流程、温度曲线和数据采集——具有重要的教学与实践价值。', Pt(14), DARK_GRAY, False),
])


# ============================================================
# SLIDE 5: Domestic vs International
# ============================================================
slide = make_content_slide(prs, '1. 研究背景及意义 · 国内外现状', 5)

# Left column - Domestic
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(0.45), MEDIUM_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(0.45),
             '国内现状', font_size=Pt(16), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.6), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('国内在建筑材料不燃性测试领域的研究主要集中在标准制定和手动测试方法的优化上。国家标准化管理委员会发布的GB/T 5464-2010标准，为实验提供了基础框架。', Pt(11), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('• 自动化程度较低，依赖人工操作', Pt(11), DARK_GRAY, False),
    ('• 软件集成度不足，缺乏统一数据管理平台', Pt(11), DARK_GRAY, False),
    ('• 高校教学场景缺少仿真实验环境', Pt(11), DARK_GRAY, False),
    ('• 硬件设备昂贵，难以规模化教学部署', Pt(11), DARK_GRAY, False),
])

# Right column - International
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45), MEDIUM_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45),
             '国外现状', font_size=Pt(16), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.3), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('在建筑材料防火测试领域，国外研究以欧美国家为主导，聚焦于智能化与自动化技术的深度整合。', Pt(11), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('• NIST标准化平台融合红外热成像技术', Pt(11), DARK_GRAY, False),
    ('• 数据驱动方法成为主流趋势', Pt(11), DARK_GRAY, False),
    ('• 专有方案成本高，数据兼容性存在挑战', Pt(11), DARK_GRAY, False),
    ('• 缺乏针对教学的仿真化定制开发案例', Pt(11), DARK_GRAY, False),
])

# Bottom summary
add_bg_rect(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.7), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(4.3), Inches(8.6), Inches(0.5),
             '💡 总结：国内外在自动化测试方面已有显著进展，但面向教学场景的低成本、纯软件仿真系统仍是空白领域。',
             font_size=Pt(12), font_color=DARK_BLUE, bold=True)


# ============================================================
# SLIDE 6: Research Significance
# ============================================================
slide = make_content_slide(prs, '1. 研究背景及意义 · 研究意义', 6)

# Left: problems
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(0.45), RED)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(0.45),
             '传统试验教学存在的问题', font_size=Pt(14), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.6), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('✘ 硬件设备昂贵：单套试验炉系统造价数十万元', Pt(11), DARK_GRAY, False),
    ('✘ 操作门槛高：需经过专业培训才能独立操作', Pt(11), DARK_GRAY, False),
    ('✘ 实验周期长：单次试验需3-4小时', Pt(11), DARK_GRAY, False),
    ('✘ 安全风险：750°C高温操作存在安全隐患', Pt(11), DARK_GRAY, False),
    ('✘ 教学效率低：一个班级需分批多次实验', Pt(11), DARK_GRAY, False),
    ('✘ 数据记录繁琐：手工记录温度，容易遗漏或错误', Pt(11), DARK_GRAY, False),
])

# Right: our solution
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45), GREEN)
add_text_box(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45),
             '仿真系统的价值', font_size=Pt(14), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.3), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('✔ 零硬件成本：纯软件运行，普通PC即可部署', Pt(11), DARK_GRAY, False),
    ('✔ 零安全风险：无高温操作，适合教学环境', Pt(11), DARK_GRAY, False),
    ('✔ 效率提升：仿真升温可加速，快速进入试验状态', Pt(11), DARK_GRAY, False),
    ('✔ 数据全自动：电子化采集，杜绝人为误差', Pt(11), DARK_GRAY, False),
    ('✔ 报告一键生成：自动生成Excel/PDF标准报告', Pt(11), DARK_GRAY, False),
    ('✔ 操作简单：图形化界面，降低学习门槛', Pt(11), DARK_GRAY, False),
])

# Bottom
add_bg_rect(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(0.7), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(4.3), Inches(8.6), Inches(0.5),
             '🎯 核心目标：开发一套完整的 ISO 11820 建筑材料不燃性试验仿真系统，服务于高校建材防火课程的教学与实训。',
             font_size=Pt(12), font_color=DARK_BLUE, bold=True)


# ============================================================
# SLIDE 7: Section 02 Divider
# ============================================================
make_section_slide(prs, '02', '系统需求分析与技术选型', 'REQUIREMENTS & TECHNOLOGY SELECTION')


# ============================================================
# SLIDE 8: Requirements Analysis
# ============================================================
slide = make_content_slide(prs, '2. 需求分析与技术选型 · 系统需求', 8)

add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.8), [
    ('本系统旨在构建一个完整的 ISO 11820 建筑材料不燃性试验仿真平台，需满足以下核心需求：', Pt(13), DARK_BLUE, True),
])

# Functional requirements in a structured list
reqs = [
    ('用户管理', '支持管理员/试验员双角色登录，基于角色的功能权限控制'),
    ('试验管理', '新建试验、填写样品信息与环境参数、试验状态全流程控制'),
    ('温度仿真', '5通道温度数据（炉温1/2、表面温、中心温、校准温）的物理仿真生成'),
    ('状态机控制', '严格遵循 Idle→Preparing→Ready→Recording→Complete 五状态流转'),
    ('实时显示', 'LED风格温度面板 + OxyPlot四通道温度曲线图（10分钟滚动窗口）'),
    ('数据存储', 'SQLite本地数据库存储试验记录，CSV文件存储每秒温度时序数据'),
    ('报告导出', '支持 Excel（含曲线图）和 PDF 两种格式自动生成'),
    ('设备校准', '仿真校准功能，支持炉壁9测温点均匀性校验和中心轴扫描'),
]

for i, (title, desc) in enumerate(reqs):
    y = Inches(1.8) + Inches(0.42) * i
    add_bg_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.35), ACCENT_BLUE)
    add_text_box(slide, Inches(0.75), y, Inches(1.6), Inches(0.35),
                 title, font_size=Pt(11), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(2.3), y, Inches(7.2), Inches(0.35),
                 desc, font_size=Pt(10), font_color=DARK_GRAY)


# ============================================================
# SLIDE 9: ISO 11820 Standards
# ============================================================
slide = make_content_slide(prs, '2. 需求分析与技术选型 · ISO 11820 标准要求', 9)

# Left: Standard requirements
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.3), Inches(3.5), LIGHT_GRAY)
add_multiline_text(slide, Inches(0.7), Inches(1.1), Inches(3.9), Inches(3.3), [
    ('ISO 11820-2020 核心要求', Pt(16), DARK_BLUE, True),
    ('', Pt(6), BLACK, False),
    ('📌 目标温度：750°C ± 5°C', Pt(12), DARK_GRAY, False),
    ('📌 稳定条件：温度在745-755°C范围内连续保持', Pt(12), DARK_GRAY, False),
    ('📌 温度漂移：10分钟漂移 ≤ 2°C', Pt(12), DARK_GRAY, False),
    ('📌 试验时长：标准60分钟，可自定义', Pt(12), DARK_GRAY, False),
    ('📌 数据采集：连续记录多通道温度与时间', Pt(12), DARK_GRAY, False),
    ('📌 判定指标：温升ΔT ≤ 50°C，失重率 ≤ 50%', Pt(12), DARK_GRAY, False),
    ('📌 质量测量：试验前后精确称重（克）', Pt(12), DARK_GRAY, False),
    ('📌 火焰观察：记录持续火焰出现时刻与时长', Pt(12), DARK_GRAY, False),
])

# Right: How we implement
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45), MEDIUM_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.45),
             '仿真实现方案', font_size=Pt(14), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.4), Inches(1.6), Inches(3.9), Inches(2.8), [
    ('🔧 SensorSimulator 温度仿真引擎', Pt(12), DARK_BLUE, True),
    ('  - 四阶段仿真：升温→稳定→记录→降温', Pt(11), DARK_GRAY, False),
    ('  - 每800ms更新一次，模拟真实采集频率', Pt(11), DARK_GRAY, False),
    ('  - 随机噪声叠加，模拟真实传感器波动', Pt(11), DARK_GRAY, False),
    ('', Pt(6), BLACK, False),
    ('🔧 TestMaster 状态机控制器', Pt(12), DARK_BLUE, True),
    ('  - 五状态严格流转，自动判定稳定条件', Pt(11), DARK_GRAY, False),
    ('  - 基于MathNet.Numerics线性回归计算温漂', Pt(11), DARK_GRAY, False),
    ('', Pt(6), BLACK, False),
    ('🔧 EPPlus + PDFsharp 报告生成', Pt(12), DARK_BLUE, True),
    ('  - Excel含三条Sheet（试验信息+温度数据+曲线图）', Pt(11), DARK_GRAY, False),
    ('  - PDF含试验概要+曲线图+判定结论', Pt(11), DARK_GRAY, False),
])


# ============================================================
# SLIDE 10: Technology Selection
# ============================================================
slide = make_content_slide(prs, '2. 需求分析与技术选型 · 技术栈', 10)

techs = [
    ('开发语言', 'C#', '.NET 8 平台，面向对象，类型安全，适合桌面应用开发'),
    ('UI 框架', 'Windows Forms (WinForms)', '成熟的桌面UI框架，控件丰富，开发效率高'),
    ('数据库', 'SQLite', '轻量级嵌入式数据库，零部署，通过 Microsoft.Data.Sqlite 直连操作'),
    ('图表组件', 'OxyPlot 2.x', '高性能开源图表库，支持实时刷新，适合动态温度曲线'),
    ('Excel导出', 'EPPlus 7.x', '支持创建带图表的专业Excel报告，无需安装Office'),
    ('PDF导出', 'PDFsharp + MigraDoc 6.x', '开源PDF生成库，支持中文，可嵌入图片与表格'),
    ('数值计算', 'MathNet.Numerics 5.x', '提供线性回归等统计算法，用于温漂计算'),
    ('日志记录', 'Serilog 4.x', '结构化日志框架，支持滚动文件输出，便于调试追踪'),
    ('配置管理', 'Microsoft.Extensions.Configuration', '读取appsettings.json，支持仿真参数灵活配置'),
]

for i, (category, tech, desc) in enumerate(techs):
    y = Inches(1.0) + Inches(0.48) * i
    add_bg_rect(slide, Inches(0.5), y, Inches(2.0), Inches(0.4), DARK_BLUE)
    add_text_box(slide, Inches(0.5), y + Inches(0.03), Inches(2.0), Inches(0.35),
                 category, font_size=Pt(10), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.7), y + Inches(0.03), Inches(2.5), Inches(0.35),
                 tech, font_size=Pt(11), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(5.4), y + Inches(0.03), Inches(4.1), Inches(0.35),
                 desc, font_size=Pt(10), font_color=DARK_GRAY)


# ============================================================
# SLIDE 11: Simulation Engine Algorithm
# ============================================================
slide = make_content_slide(prs, '2. 需求分析与技术选型 · 温度仿真算法设计', 11)

add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.5), [
    ('SensorSimulator 是系统的核心仿真引擎，负责生成5个通道的温度数据，模拟真实的物理升温过程：', Pt(12), DARK_GRAY, False),
])

# Algorithm phases
phases = [
    ('升温阶段 (TF1 < 747°C)', ACCENT_BLUE, [
        'TF1 += HeatingRate × 0.8 + Noise',
        'TF2 += HeatingRate × 0.8 + Noise (独立)',
        'TS = TF1 × 0.3 + Noise',
        'TC = TF1 × 0.25 + Noise',
        'TCal = TF1 + Noise × 2',
        '配置升温速率默认40°C/s，快速进入就绪状态',
    ]),
    ('稳定阶段 (TF1 ≥ 747°C)', GREEN, [
        'TF1 = 750 + Noise（直接钳位）',
        'TF2 = 750 + Noise',
        '稳定计数器累加，>3次tick(~3.2s)→IsStable=true',
        '同时满足745~755°C且IsStable→切换到Ready',
    ]),
    ('记录阶段 (Recording)', ORANGE, [
        'surfaceTarget = min(TF1×0.95, 800)',
        'TS += (surfaceTarget - TS) × 0.02 + Noise',
        'centerTarget = min(TF1×0.85, 750)',
        'TC += (centerTarget - TC) × 0.01 + Noise（比表面更慢）',
    ]),
    ('降温阶段 (停止加热后)', MED_GRAY, [
        'TF1 -= 0.5 + Noise × 0.1（缓慢冷却）',
        'TF2 -= 0.5 + Noise × 0.1',
        '模拟炉体自然散热过程',
    ]),
]

for i, (phase_name, color, details) in enumerate(phases):
    x = Inches(0.3) + Inches(2.4) * (i % 4)
    y = Inches(1.7)
    w = Inches(2.2)

    add_bg_rect(slide, x, y, w, Inches(0.4), color)
    add_text_box(slide, x, y + Inches(0.02), w, Inches(0.35),
                 phase_name, font_size=Pt(9), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    detail_text = '\n'.join(details)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.5), w - Inches(0.1), Inches(2.8),
                 detail_text, font_size=Pt(8), font_color=DARK_GRAY)

# Noise explanation at bottom
add_bg_rect(slide, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.45), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(4.65), Inches(8.6), Inches(0.35),
             '💡 随机噪声 = Random(-1, 1) × TempFluctuation（默认0.5°C），使仿真数据更接近真实传感器读数特征。',
             font_size=Pt(11), font_color=DARK_BLUE, bold=False)


# ============================================================
# SLIDE 12: Algorithm Summary
# ============================================================
slide = make_content_slide(prs, '2. 需求分析与技术选型 · 仿真算法小结', 12)

add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(3.5), [
    ('仿真引擎的设计原则', Pt(16), DARK_BLUE, True),
    ('', Pt(8), DARK_GRAY, False),
    ('1. 物理合理性：温度变化遵循热力学基本规律——升温速率可控、指数逼近稳态、降温自然衰减，确保仿真数据在物理意义上合理可信。', Pt(13), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('2. 通道差异性：5个通道各具独立的行为模型——炉温主导、表面温中速跟随、中心温慢速深入、校准温独立波动——真实还原多传感器系统的数据特征。', Pt(13), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('3. 可配置性：所有关键参数（目标温度、升温速率、噪声幅度、稳定阈值）均通过 appsettings.json 集中管理，无需修改代码即可调整仿真行为。', Pt(13), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('4. 阶段完整性：完整覆盖 Idle→升温→稳定→记录→降温 全生命周期，每个阶段的算法相互独立、接口统一，便于单元测试与独立调试。', Pt(13), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('5. 噪声真实性：在每个温度值上叠加可控幅度的随机噪声（±0.5°C），模拟真实传感器的量化误差与环境扰动，避免"过于完美"的数据显得不真实。', Pt(13), DARK_GRAY, False),
])


# ============================================================
# SLIDE 13: Section 03 Divider
# ============================================================
make_section_slide(prs, '03', '系统总体设计', 'SYSTEM DESIGN')


# ============================================================
# SLIDE 14: System Architecture
# ============================================================
slide = make_content_slide(prs, '3. 系统总体设计 · 架构设计', 14)

# Architecture layers
layers = [
    ('表示层 (UI)', ACCENT_BLUE, [
        'WinForms 窗体应用',
        '登录窗体 / 主窗体 / 新建试验窗体',
        'OxyPlot 实时温度曲线',
        'LED风格温度数值面板',
        '系统消息日志显示',
        'UI控件状态随状态机自动切换',
    ]),
    ('业务核心层 (Core)', MEDIUM_BLUE, [
        'TestMaster 试验状态机',
        'Idle→Preparing→Ready→Recording→Complete',
        '基于事件的异步通知机制',
        'DataBroadcast 数据广播',
        '试验终止条件自动判定',
        '温漂计算（MathNet线性回归）',
    ]),
    ('服务层 (Services)', GREEN, [
        'DaqWorker 数据采集服务 (800ms)',
        'SensorSimulator 温度仿真引擎',
        'CsvExportService CSV导出',
        'ExcelExportService Excel报告',
        'PdfExportService PDF报告',
        'Serilog 结构化日志服务',
    ]),
    ('数据层 (Data)', ORANGE, [
        'SQLite 本地数据库',
        'DbHelper 封装SQL操作',
        'Microsoft.Data.Sqlite 直连',
        '6张核心数据表',
        'CSV温度时序文件',
        'appsettings.json 配置管理',
    ]),
]

for i, (layer_name, color, items) in enumerate(layers):
    x = Inches(0.3) + Inches(2.4) * i
    w = Inches(2.2)

    add_bg_rect(slide, x, Inches(0.95), w, Inches(0.45), color)
    add_text_box(slide, x, Inches(0.97), w, Inches(0.4),
                 layer_name, font_size=Pt(11), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.1), Inches(1.55) + Inches(0.28) * j,
                     w - Inches(0.2), Inches(0.25),
                     f'• {item}', font_size=Pt(8), font_color=DARK_GRAY)

# Design principles
add_bg_rect(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.0), LIGHT_GRAY)
add_multiline_text(slide, Inches(0.7), Inches(4.1), Inches(8.6), Inches(0.85), [
    ('关键设计原则：', Pt(11), DARK_BLUE, True),
    ('• 上层依赖下层，下层不感知上层 —— 通过事件（Event）向上传递数据，禁止业务层直接调用UI方法', Pt(10), DARK_GRAY, False),
    ('• 所有UI更新必须在UI线程执行 —— 后台线程通过 Invoke 安全跨线程更新控件', Pt(10), DARK_GRAY, False),
    ('• AppContext 单例模式持有全局核心对象 —— 统一管理生命周期与依赖注入', Pt(10), DARK_GRAY, False),
])


# ============================================================
# SLIDE 15: Functional Requirements
# ============================================================
slide = make_content_slide(prs, '3. 系统总体设计 · 功能需求分析', 15)

# Flow diagram text
add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(3.8), [
    ('系统将完整的 ISO 11820 试验流程转化为自动化闭环仿真流程：', Pt(13), DARK_BLUE, True),
    ('', Pt(8), DARK_GRAY, False),
])

flow_steps = [
    ('① 登录认证', '选择角色（管理员/试验员），输入密码验证身份，进入主界面'),
    ('② 新建试验', '填写样品编号、名称、规格、尺寸、初始质量，设定环境温湿度与试验时长'),
    ('③ 开始升温', '仿真引擎启动，炉温从初始值上升，实时显示5通道温度与曲线'),
    ('④ 自动判定稳定', '系统自动检测温度范围745~755°C、稳定性计数器，达标后提示"可以开始记录"'),
    ('⑤ 开始记录', '进入Recording状态，计时器启动，每秒记录一行温度数据至CSV文件'),
    ('⑥ 试验结束', '标准60分钟或自定义时长到达、或满足终止条件、或手动停止'),
    ('⑦ 填写现象', '记录是否出现火焰、试验后质量、备注，自动计算失重率和温升'),
    ('⑧ 报告生成', '一键导出Excel报告（含曲线图）和PDF报告，归档至本地文件夹'),
]

for i, (step, desc) in enumerate(flow_steps):
    y = Inches(1.7) + Inches(0.42) * i
    add_bg_rect(slide, Inches(0.5), y, Inches(2.5), Inches(0.35), DARK_BLUE)
    add_text_box(slide, Inches(0.5), y + Inches(0.02), Inches(2.5), Inches(0.3),
                 step, font_size=Pt(10), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), y + Inches(0.02), Inches(6.3), Inches(0.3),
                 desc, font_size=Pt(10), font_color=DARK_GRAY)


# ============================================================
# SLIDE 16: System Function Modules
# ============================================================
slide = make_content_slide(prs, '3. 系统总体设计 · 功能模块设计', 16)

modules = [
    ('试验管理控制模块', DARK_BLUE, '以 TestMaster 状态机为核心，定义五状态流转逻辑，各状态转换依赖严格预定义条件，确保试验全程自动化与标准化。UI按钮状态与后台状态机严格同步，防止误操作。'),
    ('温度仿真引擎模块', MEDIUM_BLUE, 'SensorSimulator 以 800ms 周期生成5通道仿真温度数据，涵盖升温、稳定、记录、降温四阶段，叠加可控随机噪声模拟真实传感器行为。'),
    ('实时数据采集模块', ACCENT_BLUE, 'DaqWorker 后台服务每800ms轮询一次，读取仿真/硬件数据，更新内存传感器字典，触发 DataBroadcast 事件通知UI刷新。'),
    ('数据存储与查询模块', GREEN, 'SQLite 存储6张业务表，CSV存储每秒温度时序数据。历史记录支持按日期范围、样品编号、操作员多条件组合查询。'),
    ('报告生成模块', ORANGE, '试验完成后自动计算失重率与温升，生成含3个Sheet的Excel报告（信息表+数据表+曲线图）和PDF报告。'),
    ('设备校准模块', MED_GRAY, '仿真校准功能：炉壁9测温点均匀性校验、中心轴15点扫描、轴向/层级偏差分析，校准结果JSON存储。'),
]

for i, (mod_name, color, desc) in enumerate(modules):
    y = Inches(1.0) + Inches(0.7) * i

    add_bg_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.4), color)
    add_text_box(slide, Inches(0.5), y + Inches(0.03), Inches(2.2), Inches(0.35),
                 mod_name, font_size=Pt(10), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.9), y + Inches(0.03), Inches(6.6), Inches(0.55),
                 desc, font_size=Pt(10), font_color=DARK_GRAY)


# ============================================================
# SLIDE 17: Database Design
# ============================================================
slide = make_content_slide(prs, '3. 系统总体设计 · 数据库设计', 17)

# Database tables
tables = [
    ('operators', '操作员表', 'userid, username, pwd, usertype', '登录认证，管理员/试验员角色'),
    ('productmaster', '样品信息表', 'productid(PK), productname, specific, diameter, height', '存储待测材料基本信息'),
    ('apparatus', '设备信息表', 'apparatusid(PK), innernumber, apparatusname, checkdatef/t', '设备编号、检定日期、恒功率值'),
    ('testmaster', '试验记录表 ⭐', 'productid+testid(PK), operator, preweight, postweight, lostweight_per, deltatf, maxtf1~4, finaltf1~4...', '核心表：每次试验完整记录，温升+失重率判定'),
    ('sensors', '传感器配置表', 'sensorid(PK), sensorname, dispname, outputzero/span, outputvalue', '5个温度通道的配置与当前值'),
    ('CalibrationRecords', '校准记录表', 'Id(PK), CalibrationDate, ApparatusId, TempA1~C3, TAvg, TDev...', '炉壁9点+中心轴校准历史'),
]

for i, (tbl_name, tbl_desc, fields, note) in enumerate(tables):
    y = Inches(0.95) + Inches(0.68) * i

    add_bg_rect(slide, Inches(0.3), y, Inches(2.0), Inches(0.35), DARK_BLUE)
    add_text_box(slide, Inches(0.3), y + Inches(0.02), Inches(2.0), Inches(0.3),
                 tbl_name, font_size=Pt(9), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.5), y + Inches(0.02), Inches(1.3), Inches(0.3),
                 tbl_desc, font_size=Pt(9), font_color=DARK_BLUE, bold=True)

    add_text_box(slide, Inches(3.9), y + Inches(0.02), Inches(3.0), Inches(0.3),
                 fields, font_size=Pt(7.5), font_color=DARK_GRAY)

    add_text_box(slide, Inches(7.0), y + Inches(0.02), Inches(2.7), Inches(0.3),
                 note, font_size=Pt(8), font_color=MED_GRAY)

# Bottom note
add_bg_rect(slide, Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.25), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(5.0), Inches(8.6), Inches(0.25),
             '💡 温度时序数据不入库，以CSV文件独立存储：{BaseDirectory}\\TestData\\{ProductId}\\{TestId}\\sensor_data.csv',
             font_size=Pt(9), font_color=MED_GRAY)


# ============================================================
# SLIDE 18: Section 04 Divider
# ============================================================
make_section_slide(prs, '04', '系统详细设计与实现', 'DETAILED DESIGN & IMPLEMENTATION')


# ============================================================
# SLIDE 19: Core Functions - State Machine
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 试验状态机', 19)

add_multiline_text(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.5), [
    ('采用状态机模式管理试验全流程，定义5个核心状态，严格遵循ISO 11820标准流程：', Pt(12), DARK_GRAY, False),
])

# State machine flow
states = [
    ('Idle\n空闲', DARK_BLUE, '系统初始状态\n等待用户操作'),
    ('Preparing\n升温中', MEDIUM_BLUE, '仿真引擎运行\n炉温上升中'),
    ('Ready\n就绪', GREEN, '温度稳定达标\n可开始记录'),
    ('Recording\n记录中', ACCENT_BLUE, '计时器运行\n持续记录数据'),
    ('Complete\n完成', ORANGE, '试验结束\n等待保存记录'),
]

for i, (state_name, color, desc) in enumerate(states):
    x = Inches(0.2) + Inches(1.9) * i
    w = Inches(1.7)

    # State box
    add_bg_rect(slide, x, Inches(1.7), w, Inches(0.75), color)
    add_text_box(slide, x, Inches(1.72), w, Inches(0.7),
                 state_name, font_size=Pt(11), font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < 4:
        add_text_box(slide, x + w, Inches(1.85), Inches(0.25), Inches(0.4),
                     '→', font_size=Pt(20), font_color=ACCENT_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x, Inches(2.6), w, Inches(0.6),
                 desc, font_size=Pt(8), font_color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

# Transition rules
add_bg_rect(slide, Inches(0.5), Inches(3.3), Inches(9.0), Inches(1.8), LIGHT_GRAY)
add_multiline_text(slide, Inches(0.7), Inches(3.4), Inches(8.6), Inches(1.6), [
    ('状态转换规则：', Pt(13), DARK_BLUE, True),
    ('', Pt(4), DARK_GRAY, False),
    ('Idle → Preparing：用户点击「开始升温」，仿真引擎启动', Pt(11), DARK_GRAY, False),
    ('Preparing → Ready：TF1达到745~755°C 且 稳定计数器>3次tick（约3.2秒），自动判定', Pt(11), DARK_GRAY, False),
    ('Ready → Preparing：温度跌出稳定范围时自动回退（容错机制）', Pt(11), DARK_GRAY, False),
    ('Ready → Recording：用户点击「开始记录」，计时器启动，CSV文件开始写入', Pt(11), DARK_GRAY, False),
    ('Recording → Complete：固定时长到达 / 标准模式60分钟 / 用户手动停止 / 满足终止条件', Pt(11), DARK_GRAY, False),
    ('Complete：UI提示保存试验记录（试验后质量+现象），保存后 Flag="10000000"', Pt(11), DARK_GRAY, False),
    ('⚠ 关键保护：已完成但未保存的试验（totaltesttime>0且flag≠"10000000"）禁止新建试验或重新开始记录', Pt(11), RED, True),
])


# ============================================================
# SLIDE 20: Core Functions - Data Acquisition & Simulation
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 数据采集与仿真', 20)

# Left: DaqWorker
add_bg_rect(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.45), DARK_BLUE)
add_text_box(slide, Inches(0.3), Inches(1.02), Inches(4.5), Inches(0.4),
             'DaqWorker 数据采集服务（800ms周期）', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('后台独立线程，每800ms执行一次循环：', Pt(10), DARK_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('1. 读取 appsettings.json → EnableSimulation', Pt(10), DARK_GRAY, False),
    ('2. 仿真模式 → SensorSimulator.Update() 生成温度', Pt(10), DARK_GRAY, False),
    ('3. 硬件模式 → 通过串口Modbus读取真实传感器', Pt(10), DARK_GRAY, False),
    ('4. 更新内存传感器字典（5通道数值）', Pt(10), DARK_GRAY, False),
    ('5. 触发 DataBroadcast 事件广播数据', Pt(10), DARK_GRAY, False),
    ('6. UI线程通过Invoke安全更新控件', Pt(10), DARK_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('⚡ 当前系统为纯仿真模式（EnableSimulation=true）', Pt(10), ORANGE, True),
])

# Right: CSV & Data flow
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.45), MEDIUM_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.02), Inches(4.5), Inches(0.4),
             '数据存储与流转', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.4), Inches(1.6), Inches(4.1), Inches(2.5), [
    ('实时数据（800ms更新）：', Pt(10), DARK_BLUE, True),
    ('  → 传感器字典（内存）', Pt(10), DARK_GRAY, False),
    ('  → DataBroadcast事件 → UI刷新', Pt(10), DARK_GRAY, False),
    ('  → OxyPlot曲线图实时追加数据点', Pt(10), DARK_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('记录阶段（每秒写入）：', Pt(10), DARK_BLUE, True),
    ('  → CSV文件逐行追加温度数据', Pt(10), DARK_GRAY, False),
    ('  → 路径：TestData/{ProductId}/{TestId}/', Pt(10), DARK_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('试验完成（最终落库）：', Pt(10), DARK_BLUE, True),
    ('  → testmaster表UPDATE统计字段', Pt(10), DARK_GRAY, False),
    ('  → 自动生成Excel+PDF报告', Pt(10), DARK_GRAY, False),
    ('  → 报告路径：Reports/{TestId}_报告.xlsx', Pt(10), DARK_GRAY, False),
])


# ============================================================
# SLIDE 21: Core Functions - UI & Interaction
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 用户界面与交互', 21)

features = [
    ('实时温度显示', '5通道LED风格大字体数值面板（°C，1位小数），颜色区分各通道，每秒刷新'),
    ('温度曲线图', 'OxyPlot控件4条折线（炉温1/2、表面温、中心温），X轴10分钟滚动窗口，Y轴0~800°C，不同颜色区分'),
    ('计时器', 'Recording状态下显示已记录秒数，精确到秒，非记录状态显示"--"'),
    ('温度漂移', '基于MathNet.Numerics线性回归，计算最近10分钟炉温变化趋势（°C/10min）'),
    ('状态指示', '当前状态中文显示（空闲/升温中/就绪/记录中/完成），不同状态不同背景色'),
    ('系统消息', '时间戳+事件内容，不同事件不同颜色（白色普通/黄色提示/红色警告），自动滚动'),
    ('按钮联锁', 'UI按钮可用状态严格跟随状态机，防止用户在错误时机操作（如未Ready不能开始记录）'),
    ('仿真模式', 'Edit→Parameters可调整升温速率、目标温度、噪声幅度等仿真参数，无需重启'),
]

for i, (feat_name, feat_desc) in enumerate(features):
    y = Inches(1.0) + Inches(0.52) * i
    add_bg_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.4), ACCENT_BLUE)
    add_text_box(slide, Inches(0.75), y + Inches(0.05), Inches(2.0), Inches(0.35),
                 feat_name, font_size=Pt(10), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(6.7), Inches(0.45),
                 feat_desc, font_size=Pt(9), font_color=DARK_GRAY)


# ============================================================
# SLIDE 22: Core Functions - Calibration & Export
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 设备校准与报告导出', 22)

# Left: Calibration
add_bg_rect(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.45), DARK_BLUE)
add_text_box(slide, Inches(0.3), Inches(1.02), Inches(4.5), Inches(0.4),
             '📐 设备校准系统', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.5), Inches(1.6), Inches(4.1), Inches(2.8), [
    ('炉壁温度均匀性校验：', Pt(11), DARK_BLUE, True),
    ('• 9个测温点：A/B/C层 × 1/2/3轴', Pt(10), DARK_GRAY, False),
    ('• 计算轴向均值（Axis1-3）和层级均值（LevelA-C）', Pt(10), DARK_GRAY, False),
    ('• 最大偏差判定均匀性是否达标', Pt(10), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('中心轴扫描校验：', Pt(11), DARK_BLUE, True),
    ('• 探头以10mm间隔扫描5-145mm范围', Pt(10), DARK_GRAY, False),
    ('• 记录15点温度数据', Pt(10), DARK_GRAY, False),
    ('• 评估轴向温度分布均匀性', Pt(10), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('校准数据以JSON格式存储中心轴数据，关系型字段存储炉壁校验结果，支持历史查询与报告导出。', Pt(10), MED_GRAY, False),
])

# Right: Export
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.45), GREEN)
add_text_box(slide, Inches(5.2), Inches(1.02), Inches(4.5), Inches(0.4),
             '📄 报告导出系统', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.4), Inches(1.6), Inches(4.1), Inches(2.8), [
    ('Excel 报告（EPPlus生成）：', Pt(11), GREEN, True),
    ('• Sheet1：试验信息表（样品、环境、结果）', Pt(10), DARK_GRAY, False),
    ('• Sheet2：温度数据表（逐秒全部通道）', Pt(10), DARK_GRAY, False),
    ('• Sheet3：温度曲线图（嵌入式图表）', Pt(10), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('PDF 报告（PDFsharp + MigraDoc）：', Pt(11), GREEN, True),
    ('• 试验概要信息', Pt(10), DARK_GRAY, False),
    ('• 温度曲线图片（从OxyPlot导出）', Pt(10), DARK_GRAY, False),
    ('• 判定结论（通过/不通过）', Pt(10), DARK_GRAY, False),
    ('', Pt(6), DARK_GRAY, False),
    ('CSV 温度数据：试验完成自动生成，每秒一行，含时间和5通道温度，路径按样品/试验编号层级组织。', Pt(10), MED_GRAY, False),
])

# Bottom
add_bg_rect(slide, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.35), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(4.72), Inches(8.6), Inches(0.3),
             '判定结论：ΔT（温升）≤ 50°C 且 失重率 ≤ 50% 且 火焰持续时间 < 5秒 → 判定"通过"',
             font_size=Pt(11), font_color=DARK_BLUE, bold=True)


# ============================================================
# SLIDE 23: Interface Showcase - Main
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 系统界面展示（一）', 23)

# Placeholder for screenshots
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.5), Inches(1.8), LIGHT_GRAY)
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.4), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.02), Inches(4.5), Inches(0.35),
             '系统试验主界面', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.7), Inches(1.55), Inches(4.1), Inches(1.1), [
    ('[ 系统运行截图 - 主界面 ]', Pt(14), MED_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('包含：5通道温度LED面板、OxyPlot四曲线图、', Pt(10), MED_GRAY, False),
    ('状态指示、计时器、系统消息日志、操作按钮区', Pt(10), MED_GRAY, False),
])

add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(1.8), LIGHT_GRAY)
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.4), MEDIUM_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.02), Inches(4.5), Inches(0.35),
             '新建试验窗口', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.4), Inches(1.55), Inches(4.1), Inches(1.1), [
    ('[ 系统运行截图 - 新建试验 ]', Pt(14), MED_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('包含：样品信息、环境参数、试验时长模式、', Pt(10), MED_GRAY, False),
    ('设备信息自动带入、初始质量填写', Pt(10), MED_GRAY, False),
])

# Bottom: Interface features summary
add_bg_rect(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(0.4), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(3.12), Inches(9.0), Inches(0.35),
             '界面设计特点', font_size=Pt(13), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

interface_features = [
    ('布局分区', '上方温度面板 → 中部曲线图 → 右侧控制按钮 → 底部消息日志，信息层级清晰'),
    ('颜色语义', '炉温1红色、炉温2蓝色、表面温橙色、中心温绿色，通道一目了然'),
    ('字体设计', '温度数值LED风格大字体（28pt），强调数据可读性'),
    ('状态可见', '当前状态标签颜色跟随状态变化（灰色→蓝色→绿色→橙色），直观感知试验进度'),
]

for i, (feat, desc) in enumerate(interface_features):
    y = Inches(3.65) + Inches(0.38) * i
    add_text_box(slide, Inches(0.7), y, Inches(1.5), Inches(0.3),
                 feat, font_size=Pt(10), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(2.3), y, Inches(7.2), Inches(0.3),
                 desc, font_size=Pt(10), font_color=DARK_GRAY)


# ============================================================
# SLIDE 24: Interface Showcase - Calibration & Simulation
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 系统界面展示（二）', 24)

add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.5), Inches(1.8), LIGHT_GRAY)
add_bg_rect(slide, Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.4), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.02), Inches(4.5), Inches(0.35),
             '设备校准界面', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(0.7), Inches(1.55), Inches(4.1), Inches(1.1), [
    ('[ 系统运行截图 - 校准界面 ]', Pt(14), MED_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('包含：炉壁9测温点数据、轴向/层级偏差分析、', Pt(10), MED_GRAY, False),
    ('中心轴扫描数据、校准历史记录列表', Pt(10), MED_GRAY, False),
])

add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(1.8), LIGHT_GRAY)
add_bg_rect(slide, Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.4), MEDIUM_BLUE)
add_text_box(slide, Inches(5.2), Inches(1.02), Inches(4.5), Inches(0.35),
             '仿真温度曲线示意图', font_size=Pt(12), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(5.4), Inches(1.55), Inches(4.1), Inches(1.1), [
    ('[ OxyPlot 实时温度曲线截图 ]', Pt(14), MED_GRAY, False),
    ('', Pt(4), DARK_GRAY, False),
    ('显示4条曲线：炉温1(红)、炉温2(蓝)、', Pt(10), MED_GRAY, False),
    ('表面温(橙)、中心温(绿)，X轴10分钟滚动', Pt(10), MED_GRAY, False),
])

# Simulation features
add_bg_rect(slide, Inches(0.5), Inches(3.1), Inches(9.0), Inches(0.4), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(3.12), Inches(9.0), Inches(0.35),
             '仿真引擎可配置参数', font_size=Pt(13), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

params = [
    ('初始炉温', 'InitialFurnaceTemp: 720°C', '接近目标温度，快速进入Ready'),
    ('目标温度', 'TargetFurnaceTemp: 750°C', 'ISO 11820标准要求'),
    ('升温速率', 'HeatingRatePerSecond: 40°C/s', '可调整，演示用高速，教学用低速'),
    ('温度噪声', 'TempFluctuation: 0.5°C', '模拟真实传感器波动幅度'),
    ('稳定阈值', 'StableThreshold: 3.0°C', '745~755°C范围判定'),
    ('仿真开关', 'EnableSimulation: true', '纯仿真/硬件混合模式切换'),
]

for i, (param, value, note) in enumerate(params):
    y = Inches(3.65) + Inches(0.28) * i
    add_text_box(slide, Inches(0.7), y, Inches(2.0), Inches(0.25),
                 param, font_size=Pt(9), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(3.5), Inches(0.25),
                 value, font_size=Pt(9), font_color=DARK_GRAY)
    add_text_box(slide, Inches(6.4), y, Inches(3.1), Inches(0.25),
                 note, font_size=Pt(8), font_color=MED_GRAY)


# ============================================================
# SLIDE 25: Summary & Outlook
# ============================================================
slide = make_content_slide(prs, '4. 系统详细设计与实现 · 总结与展望', 25)

# Achievements
add_bg_rect(slide, Inches(0.3), Inches(1.0), Inches(9.4), Inches(0.45), DARK_BLUE)
add_text_box(slide, Inches(0.3), Inches(1.02), Inches(9.4), Inches(0.4),
             '✅ 已完成的工作', font_size=Pt(14), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

achievements = [
    ('📋 完整仿真系统', '基于 .NET 8 + WinForms 构建了 ISO 11820 建筑材料不燃性试验全流程仿真平台，涵盖试验创建、升温控制、自动判定、数据记录、报告生成等完整功能闭环。'),
    ('🔬 五通道温度仿真引擎', 'SensorSimulator 实现了升温→稳定→记录→降温四阶段的物理仿真，5个通道各具独立行为模型，叠加随机噪声，数据逼真可信。'),
    ('⚙️ 严格状态机设计', 'TestMaster 实现 Idle→Preparing→Ready→Recording→Complete 五状态流转，严格遵循 ISO 11820 标准流程，按钮状态与后台状态机同步联锁。'),
    ('📊 多格式报告导出', '支持 CSV（自动生成）、Excel（EPPlus含曲线图）、PDF（PDFsharp-MigraDoc）三种格式，满足不同场景需求。'),
    ('💾 轻量级数据架构', 'SQLite 本地数据库 + CSV文件混合存储，零部署依赖，适合教学环境快速分发。'),
]

for i, (title, desc) in enumerate(achievements):
    y = Inches(1.6) + Inches(0.48) * i
    add_text_box(slide, Inches(0.6), y, Inches(2.2), Inches(0.4),
                 title, font_size=Pt(11), font_color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(2.9), y, Inches(6.6), Inches(0.55),
                 desc, font_size=Pt(9), font_color=DARK_GRAY)

# Future work
add_bg_rect(slide, Inches(0.3), Inches(4.0), Inches(9.4), Inches(1.1), LIGHT_GRAY)
add_multiline_text(slide, Inches(0.5), Inches(4.05), Inches(9.0), Inches(1.0), [
    ('🔮 未来工作展望', Pt(13), DARK_BLUE, True),
    ('• 架构演进：向 Web 化发展，采用 ASP.NET Core + Blazor 构建跨平台监控界面，实现多设备集中远程管理', Pt(10), DARK_GRAY, False),
    ('• 硬件对接：接入真实 Modbus RTU 硬件模块（Adam 4018+），实现从纯仿真到虚实结合的升级', Pt(10), DARK_GRAY, False),
    ('• 智能升级：引入机器学习模型，基于历史数据预测材料不燃性能，辅助判定决策', Pt(10), DARK_GRAY, False),
    ('• 平台扩展：构建模块化架构，支持更多建材防火测试标准（如 ISO 1716 热值测试），打造通用型教学平台', Pt(10), DARK_GRAY, False),
])


# ============================================================
# SLIDE 26: Thank You
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# Dark blue background
add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)

# Decorative accent bar at bottom
add_bg_rect(slide, 0, Inches(4.8), SLIDE_W, Inches(0.05), ACCENT_BLUE)

# Thank you text
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
             '敬请各位老师批评指正',
             font_size=Pt(38), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Decorative line
add_bg_rect(slide, Inches(3.5), Inches(2.7), Inches(3), Inches(0.03), ACCENT_BLUE)

# Project info
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.5),
             'ISO 11820 建筑材料不燃性试验仿真系统',
             font_size=Pt(18), font_color=RGBColor(0xCC, 0xDD, 0xFF), bold=False,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.4),
             '基于 .NET 8 + WinForms + SQLite | 纯仿真 · 零硬件 · 全流程',
             font_size=Pt(12), font_color=RGBColor(0x99, 0xAA, 0xCC), bold=False,
             alignment=PP_ALIGN.CENTER)

# Bottom bar
add_bg_rect(slide, 0, Inches(5.2), SLIDE_W, Inches(0.08), ACCENT_BLUE)


# ============================================================
# SAVE
# ============================================================
output_path = r'D:\jianzhucailiao\ISO11820_仿真系统_答辩.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
